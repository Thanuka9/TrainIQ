from flask import (
    Blueprint, request, jsonify, render_template, url_for, make_response,
    session, redirect, flash, Response, stream_with_context, current_app,
)
from dotenv import load_dotenv
from flask_login import login_required, current_user
import json as json_lib
from werkzeug.utils import secure_filename
import logging
from bson.objectid import ObjectId
from models import db, StudyMaterial, SubTopic, UserProgress, User, Level, Area, UserLevelProgress, Designation, Category, LevelArea, CourseNote
from datetime import datetime
from io import BytesIO
import PyPDF2
from docx import Document
from pptx import Presentation
from utils.level_access import (
    advance_user_level_after_completion,
    can_access_level_number,
    can_access_study_material,
)
from utils.tenant_utils import (
    assert_tenant_access,
    filter_by_user_tenant,
    tenant_categories_query,
    tenant_designations_query,
    tenant_levels_query,
    user_tenant_id,
)
from utils.mongo_tenant import get_tenant_gridfs, open_grid_file
from utils.api_errors import handle_api_exception
from utils.user_access import can_upload_study_materials, effective_is_super_admin


# Load environment variables
load_dotenv()

# Configure logging
logging.basicConfig(level=logging.INFO)

# Initialize Blueprint
study_material_routes = Blueprint('study_material_routes', __name__)

# Allowed file extensions (re-export shared constants)
from utils.media_upload_constants import (
    ALLOWED_DOC_EXTENSIONS as ALLOWED_EXTENSIONS,
    ALLOWED_VIDEO_EXTENSIONS,
    MAX_DOC_SIZE_MB as MAX_FILE_SIZE_MB,
    MAX_VIDEO_SIZE_MB,
    allowed_document as allowed_file,
    allowed_video as allowed_video_file,
)

def validate_file_size(file, max_size_mb):
    """Validate if a file's size is within the specified limit."""
    file.seek(0, 2)  # move to end
    size = file.tell()
    file.seek(0)     # reset pointer
    return size <= max_size_mb * 1024 * 1024


def _upload_form_context():
    """Server-rendered dropdown data for upload study form."""
    levels = tenant_levels_query().order_by(Level.level_number.asc()).all()
    categories = tenant_categories_query().order_by(Category.name.asc()).all()
    designations = tenant_designations_query().order_by(Designation.starting_level.asc()).all()
    return {
        "levels": levels,
        "categories": categories,
        "designations": designations,
    }


def _render_upload_study():
    return render_template("upload_study.html", **_upload_form_context())

def calculate_total_pages(file_like, filetype):
    """
    Calculate total pages/slides for a given file-like object based on its type.
    file_like should be a BytesIO or similar that we can seek(0).
    """
    try:
        file_like.seek(0)
        if filetype == 'pdf':
            reader = PyPDF2.PdfReader(file_like)
            return len(reader.pages)
        elif filetype == 'docx':
            doc = Document(file_like)
            word_count = sum(len(p.text.split()) for p in doc.paragraphs)
            # approximate pages by word count / 300
            return max(1, word_count // 300)
        elif filetype == 'pptx':
            presentation = Presentation(file_like)
            return len(presentation.slides)
        else:
            return 0  # unsupported
    except Exception as e:
        logging.error(f"Error calculating total pages for {filetype}: {e}")
        return 0

@study_material_routes.route('/upload_course', methods=['GET', 'POST'])
@login_required
def upload_course():
    """
    Handle uploading of study materials and subtopics,
    with metadata in PostgreSQL and files in MongoDB (GridFS).
    """
    try:
        # Permission check
        if not can_upload_study_materials(current_user):
            flash("You do not have permission to upload study materials.", "error")
            return redirect(url_for('study_material_routes.list_study_materials'))

        # Render form
        if request.method == 'GET':
            return _render_upload_study()

        # -------------------------
        # 1) Get Form Fields
        # -------------------------
        title       = request.form.get('title')
        description = request.form.get('description')
        course_time = request.form.get('course_time')
        max_time    = request.form.get('max_time')
        level_id    = request.form.get('level_id')
        category_id = request.form.get('category_id')
        sync_minimum = request.form.get('sync_minimum_level', '1') != '0'
        try:
            minimum_level = int(request.form.get('minimum_level') or 1)
        except ValueError:
            minimum_level = 1

        subtopic_titles = request.form.getlist('subtopic_titles')
        subtopic_files  = request.files.getlist('subtopic_files')
        link_urls       = request.form.getlist('media_link_urls')
        link_titles     = request.form.getlist('media_link_titles')
        link_transcripts = request.form.getlist('media_link_transcripts')
        video_transcripts = request.form.getlist('video_transcripts')

        # Basic validation
        if not (title and description and course_time and max_time):
            flash("All fields are required.", "error")
            return redirect(url_for('study_material_routes.upload_course'))

        main_docs = request.files.getlist('main_documents')
        main_videos = request.files.getlist('video_files')
        has_docs = any(f and f.filename for f in main_docs)
        has_videos = any(f and f.filename for f in main_videos)
        has_links = any((u or '').strip() for u in link_urls)
        if not (has_docs or has_videos or has_links):
            flash("Add at least one document, video upload, or external media link.", "error")
            return redirect(url_for('study_material_routes.upload_course'))

        from models import Tenant
        from utils.tenant_storage import assert_storage_allowed, invalidate_tenant_storage_cache, sum_upload_file_sizes

        upload_tenant = Tenant.query.get(user_tenant_id())
        pending_bytes = (
            sum_upload_file_sizes(main_docs)
            + sum_upload_file_sizes(main_videos)
            + sum_upload_file_sizes(subtopic_files)
        )
        if upload_tenant and pending_bytes and not assert_storage_allowed(upload_tenant, pending_bytes):
            return redirect(url_for('study_material_routes.upload_course'))

        try:
            course_time = int(course_time)
            max_time    = int(max_time)
        except ValueError:
            flash("Course time and max time must be integers.", "error")
            return redirect(url_for('study_material_routes.upload_course'))

        # Convert optional fk's
        try:
            level_id = int(level_id) if level_id else None
        except ValueError:
            level_id = None

        if sync_minimum and level_id:
            lvl = Level.query.get(level_id)
            if lvl and lvl.level_number:
                minimum_level = int(lvl.level_number)

        try:
            category_id = int(category_id) if category_id else None
        except ValueError:
            category_id = None

        # -------------------------
        # 2) Create StudyMaterial
        # -------------------------
        study_material = StudyMaterial(
            title=title,
            description=description,
            course_time=course_time,
            max_time=max_time,
            total_pages=0,
            files=[],
            media_assets=[],
            level_id=level_id,
            category_id=category_id,
            minimum_level=minimum_level,
            tenant_id=user_tenant_id(),
        )
        db.session.add(study_material)
        db.session.commit()
        logging.info(f"Created study material with ID: {study_material.id}")

        # -------------------------
        # 3) Main Documents + Videos + Links
        # -------------------------
        from utils.course_assets import build_document_asset, build_video_asset, build_link_asset

        files = main_docs
        video_files = main_videos
        file_ids = []
        media_assets = []
        total_pages = 0

        for file in files:
            if not (file and allowed_file(file.filename)):
                continue

            if not validate_file_size(file, MAX_FILE_SIZE_MB):
                flash(f"{file.filename} exceeds the {MAX_FILE_SIZE_MB}MB limit.", "error")
                continue

            data = file.read()
            tid = study_material.tenant_id or user_tenant_id()
            gfs = get_tenant_gridfs(tid)
            mongo_id = gfs.put(
                data,
                filename=secure_filename(file.filename),
                metadata={"tenant_id": tid, "study_material_id": study_material.id, "asset_type": "document"},
            )
            file_ids.append(f"{mongo_id}|{file.filename}")
            media_assets.append(build_document_asset(str(mongo_id), file.filename))

            pages = calculate_total_pages(BytesIO(data), file.filename.rsplit('.',1)[1].lower())
            total_pages += pages

        transcript_idx = 0
        for file in video_files:
            if not (file and allowed_video_file(file.filename)):
                continue
            if not validate_file_size(file, MAX_VIDEO_SIZE_MB):
                flash(f"{file.filename} exceeds the {MAX_VIDEO_SIZE_MB}MB video limit.", "error")
                continue

            data = file.read()
            tid = study_material.tenant_id or user_tenant_id()
            gfs = get_tenant_gridfs(tid)
            mongo_id = gfs.put(
                data,
                filename=secure_filename(file.filename),
                metadata={"tenant_id": tid, "study_material_id": study_material.id, "asset_type": "video"},
            )
            vid_id = str(mongo_id)
            file_ids.append(f"{vid_id}|{file.filename}")
            transcript = video_transcripts[transcript_idx] if transcript_idx < len(video_transcripts) else ""
            transcript_idx += 1
            media_assets.append(build_video_asset(vid_id, file.filename, transcript))
            total_pages += 1

        for idx, raw_url in enumerate(link_urls):
            url = (raw_url or '').strip()
            if not url:
                continue
            link_title = (link_titles[idx] if idx < len(link_titles) else '').strip()
            link_transcript = (link_transcripts[idx] if idx < len(link_transcripts) else '').strip()
            asset = build_link_asset(url, link_title, link_transcript)
            if not asset:
                flash(f"Unsupported or invalid link: {url}", "error")
                continue
            media_assets.append(asset)
            total_pages += 1

        study_material.files = file_ids
        study_material.media_assets = media_assets
        study_material.total_pages = max(total_pages, 1)

        if not media_assets:
            db.session.delete(study_material)
            db.session.commit()
            flash("No valid content was uploaded. Please add documents, videos, or supported links.", "error")
            return redirect(url_for('study_material_routes.upload_course'))

        db.session.commit()
        logging.info(f"Main documents uploaded for study material ID: {study_material.id}")

        # -------------------------
        # 4) Subtopics
        # -------------------------
        for idx, title in enumerate(subtopic_titles):
            if not title:
                continue

            file = subtopic_files[idx] if idx < len(subtopic_files) else None
            sub_pages = 0
            sub_file_id = None

            if file and allowed_file(file.filename):
                if not validate_file_size(file, MAX_FILE_SIZE_MB):
                    flash(f"{file.filename} exceeds size limit.", "error")
                    continue

                data = file.read()
                tid = study_material.tenant_id or user_tenant_id()
                gfs = get_tenant_gridfs(tid)
                mongo_id = gfs.put(
                    data,
                    filename=secure_filename(file.filename),
                    metadata={"tenant_id": tid, "study_material_id": study_material.id},
                )
                sub_file_id = str(mongo_id)

                sub_pages = calculate_total_pages(BytesIO(data), file.filename.rsplit('.',1)[1].lower())

            sub = SubTopic(
                title=title,
                study_material_id=study_material.id,
                file_id=sub_file_id,
                page_count=sub_pages
            )
            db.session.add(sub)
            study_material.total_pages += sub_pages

        db.session.commit()
        logging.info(f"Subtopics uploaded for study material ID: {study_material.id}")

        invalidate_tenant_storage_cache(study_material.tenant_id or user_tenant_id())
        flash("Study materials and subtopics uploaded successfully.", "success")
        return redirect(url_for('study_material_routes.list_study_materials'))

    except Exception as e:
        logging.error(f"Error in upload_course: {e}", exc_info=True)
        db.session.rollback()
        flash("An error occurred while uploading the course.", "error")
        return redirect(url_for('study_material_routes.upload_course'))

    
@study_material_routes.route('/start_course/<int:course_id>', methods=['POST'])
@login_required
def start_course(course_id):
    """
    Start a course for a user and record the start date.
    """
    try:
        user_id = session.get('user_id')
        if not user_id:
            return jsonify({'error': 'User ID is required to start the course'}), 400

        user = User.query.get(user_id)
        if not user:
            return jsonify({'error': 'User not found'}), 404

        study_material = StudyMaterial.query.get_or_404(course_id)
        assert_tenant_access(study_material)

        # Check access eligibility via our unified helper
        if not can_access_study_material(user, study_material):
            return jsonify({
                'error': 'You must complete earlier levels to access this course.'
            }), 403

        # Check if already started
        user_progress = UserProgress.query.filter_by(
            user_id=user_id,
            study_material_id=course_id
        ).first()

        if user_progress:
            logging.info(f"User {user_id} has already started course {course_id}.")
            return jsonify({
                'message': 'Course already started',
                'redirect_url': url_for('study_material_routes.view_course', course_id=course_id)
            }), 200

        # Start new course progress
        user_progress = UserProgress(
            user_id=user_id,
            study_material_id=course_id,
            pages_visited=0,
            progress_percentage=0,
            completion_date=None,
            start_date=datetime.utcnow()
        )
        db.session.add(user_progress)
        db.session.commit()

        logging.info(f"User {user_id} started course {course_id} at {user_progress.start_date}.")
        return jsonify({
            'success': 'Course started successfully',
            'start_date': user_progress.start_date.isoformat(),
            'redirect_url': url_for('study_material_routes.view_course', course_id=course_id)
        }), 201

    except Exception as e:
        logging.error(f"Error starting course: {e}", exc_info=True)
        return jsonify({'error': 'Failed to start course'}), 500


def can_access_level(user, level_id):
    """Backward-compatible wrapper — ``level_id`` is a curriculum level number."""
    return can_access_level_number(user, level_id)

# ----  Course Details  ----
@study_material_routes.route("/view_course/<int:course_id>")
@login_required
def view_course(course_id):
    """
    Dashboard-style page that shows title, description,
    dates, overall progress, and a “Continue” link.
    No heavy file streaming happens here.
    """
    # 1 Fetch objects --------------------------------------------------
    study_material = StudyMaterial.query.get_or_404(course_id)
    assert_tenant_access(study_material)
    subtopics      = SubTopic.query.filter_by(study_material_id=course_id).all()

    user_id = session.get("user_id")
    if not user_id:
        flash("Please log in.", "warning")
        return redirect(url_for("auth_routes.login"))

    user = User.query.get_or_404(user_id)
    if not can_access_study_material(user, study_material):
        flash("Complete previous levels to unlock this material.", "danger")
        return redirect(url_for("study_material_routes.list_study_materials"))
    # 2 Check user progress ---------------------------------------------
    user_progress = (UserProgress.query
                     .filter_by(user_id=user_id, study_material_id=course_id)
                     .first())

    if not user_progress:
        user_progress = UserProgress(
            user_id=user_id,
            study_material_id=course_id,
            pages_visited=0,
            progress_percentage=0,
            start_date=datetime.utcnow()
        )
        db.session.add(user_progress)
        db.session.commit()

    # 2 Pick first viewer asset (documents, videos, or links — not legacy files-only)
    from utils.course_assets import (
        assets_from_study_material,
        course_viewer_mode,
        course_media_summary,
    )

    assets = assets_from_study_material(study_material)
    viewer_mode = course_viewer_mode(study_material)
    asset_count = len(assets)
    ap = (user_progress.asset_progress or {}) if user_progress else {}
    assets_completed = sum(1 for a in assets if int(ap.get(str(a.get("id")), 0)) >= 100)
    first_asset_id = str(assets[0]["id"]) if assets else None
    continue_url = (
        url_for(
            "study_material_routes.course_content",
            course_id=course_id,
            file_id=first_asset_id,
        )
        if first_asset_id
        else None
    )

    # 3 Render
    return render_template(
        "view_course.html",
        study_material=study_material,
        subtopics=subtopics,
        user_progress=user_progress,
        continue_url=continue_url,
        viewer_mode=viewer_mode,
        asset_count=asset_count,
        assets_completed=assets_completed,
        media_summary=course_media_summary(study_material),
    )

# ----  Document Viewer  --------------------------------------------
@study_material_routes.route("/course_content/<int:course_id>")
@login_required
def course_content(course_id):
    """
    Streams PDFs / other docs in a dedicated viewer.
    ?file_id=<mongo-id> tells the page which file to open first.
    """

    user_id = session.get("user_id")
    if not user_id:
        flash("Please log in.", "warning")
        return redirect(url_for("auth_routes.login"))

    study_material = StudyMaterial.query.get_or_404(course_id)
    assert_tenant_access(study_material)

    # ---- Collect assets ------------------------------------------------
    requested_id = request.args.get("file_id") or request.args.get("asset_id")
    stream_tpl = url_for("study_material_routes.stream_file", file_id="__FID__")
    from utils.course_assets import viewer_assets, course_viewer_mode

    documents = viewer_assets(study_material, stream_tpl)
    viewer_mode = course_viewer_mode(study_material)

    for doc in documents:
        if doc.get("type") == "txt" and doc.get("stream_url"):
            try:
                gfile, _ = open_grid_file(doc["mongo_id"], study_material.tenant_id)
                doc["content"] = gfile.read().decode("utf-8", errors="ignore")
            except Exception as e:
                current_app.logger.warning(f"Text fetch failed: {e}")

    if requested_id:
        documents.sort(key=lambda d: 0 if str(d.get("id")) == str(requested_id) else 1)

    # ---- Progress record (unchanged) --------------------------------
    user_progress = (UserProgress.query
                     .filter_by(user_id=user_id, study_material_id=course_id)
                     .first())

    course_notes = {
        f"{n.asset_id}:{n.page_num}": n.content
        for n in CourseNote.query.filter_by(
            user_id=user_id,
            study_material_id=course_id,
        ).all()
    }

    return render_template(
        "course_content.html",
        study_material=study_material,
        documents=documents,
        user_progress=user_progress,
        viewer_mode=viewer_mode,
        asset_progress=(user_progress.asset_progress if user_progress else {}) or {},
        course_notes=course_notes,
    )

@study_material_routes.route('/list', methods=['GET'])
@login_required
def list_study_materials():
    """
    Render the list of all study materials with progress.
    """
    user_id = session.get('user_id')
    if not user_id:
        return redirect(url_for('auth_routes.login'))

    user = User.query.get(user_id)
    if not user:
        flash("User not found.", "danger")
        return redirect(url_for('auth_routes.login'))

    # Get all study materials
    materials = filter_by_user_tenant(StudyMaterial.query, StudyMaterial).all()
    accessible_materials = []

    # Filter accessible study materials
    for material in materials:
        if can_access_study_material(user, material):
            accessible_materials.append(material)

    # Prepare progress data
    progress_data = []
    for material in accessible_materials:
        user_progress = UserProgress.query.filter_by(
            user_id=user.id,
            study_material_id=material.id
        ).first()
        progress_percentage = user_progress.progress_percentage if user_progress else 0
        progress_data.append({'course_id': material.id, 'progress_percentage': progress_percentage})

    # --- SORT by Course ID ---
    # Zip, sort, and unzip to keep both lists in sync
    combined = sorted(zip(accessible_materials, progress_data), key=lambda x: x[0].id)
    if combined:
        accessible_materials, progress_data = zip(*combined)
        accessible_materials, progress_data = list(accessible_materials), list(progress_data)
    else:
        accessible_materials, progress_data = [], []

    from utils.learner_recommendations import get_recommended_next_course

    return render_template(
        'list_study_materials.html',
        materials=accessible_materials,
        progress_data=progress_data,
        recommended_course=get_recommended_next_course(user),
    )


@study_material_routes.route('/upload_page', methods=['GET'])
@login_required
def upload_page():
    """
    Render the upload page only for authorized users.
    """
    if not can_upload_study_materials(current_user):
        logging.warning(
            "Unauthorized access attempt to upload page by user ID: %s",
            session.get('user_id'),
        )
        flash("You do not have permission to upload study materials.", "danger")
        return redirect(url_for('study_material_routes.list_study_materials'))
    return _render_upload_study()


@study_material_routes.route('/study_materials', methods=['GET'])
def study_materials():
    """
    Render the Study Materials dashboard.
    """
    return render_template('study_materials.html')


@study_material_routes.route("/update_progress", methods=["POST"])
@login_required
def update_progress():
    """
    Called by the viewer whenever a page becomes 50 % visible.
    Updates pages_visited, progress %, completion_date, and (optionally) bumps the user level.
    """
    try:
        data              = request.json or {}
        user_id           = current_user.id
        study_material_id = data.get("study_material_id")
        asset_id          = data.get("asset_id")
        watch_percent     = data.get("watch_percent")
        current_page      = int(data.get("current_page", 0))
        total_pages       = int(data.get("total_pages", 0))

        if not (user_id and study_material_id):
            return jsonify(error="invalid input"), 400

        study_material = StudyMaterial.query.get_or_404(study_material_id)
        assert_tenant_access(study_material)

        user = User.query.get(user_id)
        if not can_access_study_material(user, study_material):
            return jsonify(error="Access denied"), 403

        from utils.course_progress import recalc_user_progress

        prog = (UserProgress.query
                .filter_by(user_id=user_id, study_material_id=study_material_id)
                .with_for_update()
                .first())

        if not prog:
            prog = UserProgress(
                user_id=user_id,
                study_material_id=study_material_id,
                pages_visited=0,
                asset_progress={},
                start_date=datetime.utcnow(),
            )
            db.session.add(prog)

        # Per-asset video / embed / unified asset progress
        if asset_id is not None and watch_percent is not None:
            ap = dict(prog.asset_progress or {})
            aid = str(asset_id)
            pct = min(100, max(0, int(watch_percent)))
            ap[aid] = max(int(ap.get(aid, 0)), pct)
            prog.asset_progress = ap
            recalc_user_progress(study_material, prog)
        elif total_pages > 0 and current_page > 0:
            if total_pages != study_material.total_pages:
                total_pages = study_material.total_pages or total_pages
            if current_page > prog.pages_visited:
                prog.pages_visited = min(current_page, total_pages)
            if asset_id:
                ap = dict(prog.asset_progress or {})
                ap[str(asset_id)] = min(100, int(current_page / total_pages * 100))
                prog.asset_progress = ap
                recalc_user_progress(study_material, prog)
            else:
                raw_pct = int(prog.pages_visited / total_pages * 100)
                prog.progress_percentage = min(raw_pct, 100)
                prog.completed = prog.progress_percentage >= 100
        else:
            return jsonify(error="invalid input"), 400

        if prog.progress_percentage >= 100 and prog.completion_date is None:
            prog.progress_percentage = 100
            prog.completion_date = datetime.utcnow()
            prog.completed = True

        db.session.commit()

        # ---------- level-unlock check --------------------------------
        if study_material.level_id and prog.completed:
            unlocked = advance_user_level_after_completion(
                user_id, study_material.level_id
            )
            if unlocked:
                flash(f"🎉 Level {unlocked} unlocked!", "success")

        return jsonify(
            success=True,
            progress_percentage=prog.progress_percentage,
            completed=prog.completed
        ), 200

    except Exception as e:
        logging.exception("update_progress failed")
        return handle_api_exception(e, user_message="Could not update progress.")


@study_material_routes.route("/course_note", methods=["POST"])
@login_required
def save_course_note():
    """Upsert or delete learner notes for the current asset/page scope."""
    try:
        data = request.get_json(silent=True) or {}
        study_material_id = data.get("study_material_id")
        asset_id = str(data.get("asset_id") or "")
        page_num = int(data.get("page_num") or 0)
        content = (data.get("content") or "")[:50000]

        if not study_material_id or not asset_id:
            return jsonify({"error": "Invalid input"}), 400

        study_material = StudyMaterial.query.get_or_404(study_material_id)
        assert_tenant_access(study_material)

        note = CourseNote.query.filter_by(
            user_id=current_user.id,
            study_material_id=study_material_id,
            asset_id=asset_id,
            page_num=page_num,
        ).first()

        if not content.strip():
            if note:
                db.session.delete(note)
                db.session.commit()
            return jsonify({"success": True, "deleted": True})

        if not note:
            note = CourseNote(
                user_id=current_user.id,
                study_material_id=study_material_id,
                asset_id=asset_id,
                page_num=page_num,
                tenant_id=study_material.tenant_id,
            )
            db.session.add(note)
        note.content = content.strip()
        note.updated_at = datetime.utcnow()
        db.session.commit()
        return jsonify({
            "success": True,
            "updated_at": note.updated_at.isoformat() + "Z",
        })
    except Exception as e:
        logging.exception("save_course_note failed")
        return handle_api_exception(e, user_message="Could not save note.")


def _course_notes_for_user(course_id: int, user_id: int) -> list[CourseNote]:
    study_material = StudyMaterial.query.get_or_404(course_id)
    assert_tenant_access(study_material)
    return (
        CourseNote.query.filter_by(
            user_id=user_id,
            study_material_id=course_id,
        )
        .order_by(CourseNote.asset_id, CourseNote.page_num)
        .all()
    )


def _format_course_notes_text(notes: list[CourseNote], title: str = "Course Notes") -> str:
    lines = [title, "=" * len(title), ""]
    if not notes:
        lines.append("(No notes saved yet.)")
        return "\n".join(lines)
    for n in notes:
        lines.append(f"[Asset {n.asset_id} · Page {n.page_num}]")
        lines.append(n.content or "")
        lines.append("")
    return "\n".join(lines).strip() + "\n"


@study_material_routes.route("/course_notes/search", methods=["GET"])
@login_required
def search_course_notes():
    """Search learner notes across all courses (JSON)."""
    q = (request.args.get("q") or "").strip()
    if not q or len(q) < 2:
        return jsonify({"query": q, "notes": []})

    pattern = f"%{q}%"
    rows = (
        CourseNote.query.filter_by(user_id=current_user.id)
        .filter(CourseNote.content.ilike(pattern))
        .order_by(CourseNote.updated_at.desc())
        .limit(50)
        .all()
    )

    user_tid = user_tenant_id()
    results = []
    seen = set()
    for note in rows:
        material = StudyMaterial.query.get(note.study_material_id)
        if not material:
            continue
        if (
            user_tid is not None
            and material.tenant_id is not None
            and material.tenant_id != user_tid
        ):
            continue
        key = (note.study_material_id, note.asset_id, note.page_num)
        if key in seen:
            continue
        seen.add(key)
        snippet = (note.content or "").strip()
        if len(snippet) > 160:
            snippet = snippet[:157] + "…"
        results.append({
            "course_id": note.study_material_id,
            "course_title": material.title,
            "asset_id": note.asset_id,
            "page_num": note.page_num,
            "snippet": snippet,
            "updated_at": note.updated_at.isoformat() if note.updated_at else None,
        })

    return jsonify({"query": q, "notes": results, "count": len(results)})


@study_material_routes.route("/course_notes/export/<int:course_id>", methods=["GET"])
@login_required
def export_course_notes(course_id):
    """Download all learner notes for a course as plain text."""
    from flask import Response

    notes = _course_notes_for_user(course_id, current_user.id)
    material = StudyMaterial.query.get_or_404(course_id)
    body = _format_course_notes_text(notes, title=f"{material.title} — My Notes")
    filename = f"notes_{course_id}_{current_user.id}.txt"
    return Response(
        body,
        mimetype="text/plain; charset=utf-8",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@study_material_routes.route("/course_notes/copy/<int:course_id>", methods=["GET"])
@login_required
def copy_course_notes(course_id):
    """Return all notes as JSON for clipboard copy in the UI."""
    notes = _course_notes_for_user(course_id, current_user.id)
    material = StudyMaterial.query.get_or_404(course_id)
    return jsonify({
        "success": True,
        "course_id": course_id,
        "course_title": material.title,
        "text": _format_course_notes_text(notes, title=f"{material.title} — My Notes"),
        "count": len(notes),
    })


@study_material_routes.route('/stream_file/<file_id>', methods=['GET'])
@login_required
def stream_file(file_id):
    """Stream file content for inline display (tenant-scoped)."""
    from utils.file_access import require_material_file_access

    try:
        _material, grid_file = require_material_file_access(file_id)
        filename = grid_file.filename
        extension = f".{filename.rsplit('.', 1)[-1].lower()}"
        content_type_map = {
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.pdf': 'application/pdf',
            '.txt': 'text/plain',
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.mp4': 'video/mp4',
            '.webm': 'video/webm',
            '.mov': 'video/quicktime',
            '.m4v': 'video/x-m4v',
        }
        content_type = content_type_map.get(extension, None)

        if not content_type:
            content_type = 'application/octet-stream'

        # Stream the file content in chunks
        def generate():
            try:
                while chunk := grid_file.read(8192):  # Read in 8KB chunks
                    yield chunk
            except Exception as e:
                logging.error(f"Error reading file ID {file_id} in chunks: {e}")
                raise

        # Prepare the response with appropriate headers
        response = make_response(generate())
        response.headers['Content-Type'] = content_type
        response.headers['Content-Disposition'] = f'inline; filename="{filename}"'
        response.headers['Cache-Control'] = 'no-store'  # Prevent caching
        response.headers['Content-Security-Policy'] = "frame-ancestors 'self'; script-src 'self';"
        response.headers['X-Content-Type-Options'] = 'nosniff'
        response.headers['Referrer-Policy'] = 'no-referrer'

        logging.info(f"Successfully streamed file with ID {file_id} ({filename})")
        return response

    except FileNotFoundError:
        logging.error(f"File with ID {file_id} does not exist in GridFS.")
        return jsonify({'error': 'File not found'}), 404
    except Exception as e:
        logging.error(f"Error streaming file with ID {file_id}: {e}")
        return jsonify({'error': 'Failed to stream file.'}), 500


@study_material_routes.route('/update_time', methods=['POST'])
@login_required
def update_time():
    """
    Add elapsed seconds to UserProgress.total_time.
    The viewer sends chunks (default 30 s) while the tab is visible.
    """
    try:
        data         = request.json or {}
        delta        = int(data.get('elapsed_time', 0))
        material_id  = data.get('study_material_id')
        user_id      = session.get('user_id')

        if not user_id or not material_id:
            return jsonify(error="missing ids"), 400
        if delta <= 0:
            return jsonify(success=True)          # ignore zero/neg chunks

        prog = (UserProgress.query
                .filter_by(user_id=user_id, study_material_id=material_id)
                .with_for_update()
                .first())

        if not prog:
            return jsonify(error="progress not found"), 404

        # sanity-check: don’t let a single chunk exceed total possible time.
        if prog.start_date:
            max_allowed = (datetime.utcnow() - prog.start_date).total_seconds() + 300
            if delta > max_allowed:
                return jsonify(error="elapsed_time too large"), 400

        prog.time_spent = (prog.time_spent or 0) + delta
        db.session.commit()
        return jsonify(success=True)

    except Exception as e:
        logging.exception("update_time failed")
        return handle_api_exception(e, user_message="Could not update progress.")


@study_material_routes.route('/download_file/<file_id>', methods=['GET'])
@login_required
def download_file(file_id):
    """Download a study file (tenant-scoped)."""
    from utils.file_access import require_material_file_access

    try:
        _material, grid_file = require_material_file_access(file_id)
        if not grid_file:
            logging.error(f"File with ID {file_id} not found in GridFS.")
            return jsonify({'error': 'File not found'}), 404

        # Create response with the file data
        filename = grid_file.filename
        response = make_response(grid_file.read())
        response.headers['Content-Type'] = 'application/octet-stream'
        response.headers['Content-Disposition'] = f'attachment; filename="{filename}"'
        logging.info(f"File downloaded successfully with ID {file_id} ({filename})")
        return response
    except Exception as e:
        logging.error(f"Error downloading file with ID {file_id}: {e}")
        return jsonify({'error': 'Failed to download file.'}), 500

@study_material_routes.route('/dashboard', methods=['GET'])
def dashboard():
    """
    Render the study materials dashboard with super admin access check.
    """
    # Ensure user is logged in
    user_id = session.get('user_id')
    if not user_id:
        return redirect(url_for('auth_routes.login'))

    user = User.query.get(user_id)
    return render_template(
        'dashboard.html',
        is_super_admin=effective_is_super_admin(user) if user else False,
    )


def extract_text_from_gridfs(file_id, page_num=None, tenant_id=None):
    """Extract plain text from a GridFS document for LearnIQ context."""
    try:
        gfile, _ = open_grid_file(file_id, tenant_id)
        raw = gfile.read()
        filename = gfile.filename or ""
        ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else "txt"
        buf = BytesIO(raw)

        if ext == "pdf":
            reader = PyPDF2.PdfReader(buf)
            if page_num and 1 <= page_num <= len(reader.pages):
                idx = page_num - 1
                window = []
                for i in range(max(0, idx - 1), min(len(reader.pages), idx + 2)):
                    t = reader.pages[i].extract_text() or ""
                    if t.strip():
                        label = "CURRENT PAGE" if i == idx else f"Page {i + 1}"
                        window.append(f"[{label}]\n{t}")
                return "\n\n".join(window)
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        if ext == "docx":
            doc = Document(buf)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        if ext == "pptx":
            prs = Presentation(buf)
            slides = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slides.append(shape.text)
            return "\n".join(slides)
        if ext == "txt":
            return raw.decode("utf-8", errors="ignore")
        return ""
    except Exception as e:
        logging.error(f"Text extraction failed for {file_id}: {e}")
        return ""


def _get_course_document_text(study_material, file_id=None, page_num=None):
    """Build combined text context from course files and media assets."""
    from utils.course_assets import collect_ai_text

    if file_id:
        text = collect_ai_text(study_material, file_id=str(file_id), page_num=page_num)
        if text.strip():
            return text

    tid = study_material.tenant_id
    parts = []
    for entry in (study_material.files or []):
        if "|" not in entry:
            continue
        fid, _ = (p.strip() for p in entry.split("|", 1))
        text = extract_text_from_gridfs(fid, page_num=page_num, tenant_id=tid) if page_num else extract_text_from_gridfs(fid, tenant_id=tid)
        if text.strip():
            parts.append(text)

    asset_text = collect_ai_text(study_material)
    if asset_text.strip():
        return asset_text

    return "\n\n".join(parts)


@study_material_routes.route('/ai/status', methods=['GET'])
@login_required
def learniq_status():
    """Check if local Ollama/Gemma 4 is available."""
    from utils.local_ai import get_ai_status
    return jsonify(get_ai_status())


def _ai_unavailable_response():
    from utils.local_ai import get_ai_status
    status = get_ai_status()
    return jsonify({"error": status["message"], **status}), 503


def _learniq_context(study_material, data):
    """Parse LearnIQ request fields; allow vision fallback when text is sparse."""
    from utils.local_ai import resolve_model, needs_vision_fallback
    from utils import ai_cache
    from utils.course_assets import VIDEO_TYPES

    file_id = data.get("file_id")
    page_num = data.get("current_page")
    asset_type = (data.get("asset_type") or "").strip().lower()
    video_time = data.get("video_time_seconds")

    doc_text = _get_course_document_text(study_material, file_id, page_num=page_num)

    if asset_type in VIDEO_TYPES:
        extra_bits = []
        if video_time is not None:
            try:
                sec = int(float(video_time))
                extra_bits.append(f"Learner is at approximately {sec} seconds into this video.")
            except (TypeError, ValueError):
                pass
        if not doc_text.strip() and study_material.description:
            doc_text = study_material.description.strip()
        if doc_text.strip() and "No transcript" in doc_text:
            extra_bits.append(
                "No transcript was uploaded for this video. Use title, URL, and course context; "
                "clearly note when the summary is limited."
            )
        if extra_bits:
            doc_text = (doc_text + "\n\n" + "\n".join(extra_bits)).strip()

    page_image = None
    is_video_asset = asset_type in VIDEO_TYPES
    if not is_video_asset:
        if data.get("use_vision") and data.get("page_image"):
            page_image = data.get("page_image")
        elif not doc_text.strip() and data.get("page_image"):
            page_image = data.get("page_image")
    page_images = [page_image] if page_image else None

    if not doc_text.strip() and not page_images:
        return None, jsonify({
            "error": (
                "No content for AI on this item. Add a video transcript when uploading, "
                "or open a document page with extractable text."
            )
        }), 400

    model = resolve_model()
    cache_key = ai_cache.make_key(
        data.get("_cache_feature", "learniq"),
        study_material.id,
        file_id,
        page_num,
        asset_type,
        video_time,
        model,
    )
    return {
        "file_id": file_id,
        "page_num": page_num,
        "asset_type": asset_type,
        "video_time_seconds": video_time,
        "page_images": page_images,
        "doc_text": doc_text,
        "model": model,
        "cache_key": cache_key,
        "use_vision": needs_vision_fallback(doc_text, page_images),
    }, None, None


@study_material_routes.route('/ai/summarize/<int:course_id>', methods=['POST'])
@login_required
def learniq_summarize_route(course_id):
    from utils.local_ai import learniq_summarize, is_available
    from utils.ai_rate_limit import check_ai_rate_limit

    if not session.get("user_id"):
        return jsonify({"error": "Unauthorized"}), 401
    ok, retry = check_ai_rate_limit()
    if not ok:
        return jsonify({"error": f"Rate limit exceeded. Retry in {retry}s.", "retry_after": retry}), 429

    study_material = StudyMaterial.query.get_or_404(course_id)
    assert_tenant_access(study_material)
    user = User.query.get(session["user_id"])
    if not can_access_study_material(user, study_material):
        return jsonify({"error": "Access denied"}), 403
    if not is_available():
        return _ai_unavailable_response()

    data = request.get_json(silent=True) or {}
    data["_cache_feature"] = "summarize"
    ctx, err_resp, err_code = _learniq_context(study_material, data)
    if err_resp:
        return err_resp, err_code

    from utils import ai_cache
    cached = ai_cache.get(ctx["cache_key"])
    if cached:
        return jsonify({"summary": cached["summary"], "feature": "LearnIQ", "cached": True})

    try:
        summary = learniq_summarize(
            study_material.title, ctx["doc_text"], page_images=ctx["page_images"]
        )
        ai_cache.set(ctx["cache_key"], {"summary": summary})
        return jsonify({
            "summary": summary,
            "feature": "LearnIQ",
            "cached": False,
            "vision": ctx["use_vision"],
        })
    except ConnectionError as e:
        return jsonify({"error": "AI service is temporarily unavailable."}), 503
    except Exception as e:
        logging.exception("LearnIQ summarize failed")
        return jsonify({"error": str(e) or "AI request failed."}), 500


@study_material_routes.route('/ai/stream/summarize/<int:course_id>', methods=['POST'])
@login_required
def learniq_stream_summarize_route(course_id):
    from utils.local_ai import learniq_summarize_stream, is_available
    from utils.ai_rate_limit import check_ai_rate_limit

    if not session.get("user_id"):
        return jsonify({"error": "Unauthorized"}), 401
    ok, retry = check_ai_rate_limit()
    if not ok:
        return jsonify({"error": f"Rate limit exceeded. Retry in {retry}s.", "retry_after": retry}), 429

    study_material = StudyMaterial.query.get_or_404(course_id)
    assert_tenant_access(study_material)
    user = User.query.get(session["user_id"])
    if not can_access_study_material(user, study_material):
        return jsonify({"error": "Access denied"}), 403
    if not is_available():
        return _ai_unavailable_response()

    data = request.get_json(silent=True) or {}
    data["_cache_feature"] = "summarize"
    ctx, err_resp, err_code = _learniq_context(study_material, data)
    if err_resp:
        return err_resp, err_code

    from utils import ai_cache
    cached = ai_cache.get(ctx["cache_key"])

    def generate():
        if cached:
            yield f"data: {json_lib.dumps({'text': cached['summary'], 'cached': True})}\n\n"
            yield f"data: {json_lib.dumps({'done': True, 'cached': True})}\n\n"
            return
        parts = []
        try:
            for chunk in learniq_summarize_stream(
                study_material.title, ctx["doc_text"], page_images=ctx["page_images"]
            ):
                parts.append(chunk)
                yield f"data: {json_lib.dumps({'text': chunk})}\n\n"
            summary = "".join(parts)
            ai_cache.set(ctx["cache_key"], {"summary": summary})
            yield f"data: {json_lib.dumps({'done': True, 'cached': False, 'vision': ctx['use_vision']})}\n\n"
        except ConnectionError as e:
            logging.error("LearnIQ summarize stream failed: %s", e)
            yield f"data: {json_lib.dumps({'error': 'AI service unavailable. Check that Ollama is running and try again.'})}\n\n"
        except Exception as e:
            logging.exception("LearnIQ summarize stream error")
            yield f"data: {json_lib.dumps({'error': str(e) or 'AI request failed.'})}\n\n"

    return Response(stream_with_context(generate()), mimetype="text/event-stream")


@study_material_routes.route('/ai/flashcards/<int:course_id>', methods=['POST'])
@login_required
def learniq_flashcards_route(course_id):
    from utils.local_ai import learniq_flashcards, is_available
    from utils.ai_rate_limit import check_ai_rate_limit

    if not session.get("user_id"):
        return jsonify({"error": "Unauthorized"}), 401
    ok, retry = check_ai_rate_limit()
    if not ok:
        return jsonify({"error": f"Rate limit exceeded. Retry in {retry}s.", "retry_after": retry}), 429

    study_material = StudyMaterial.query.get_or_404(course_id)
    assert_tenant_access(study_material)
    user = User.query.get(session["user_id"])
    if not can_access_study_material(user, study_material):
        return jsonify({"error": "Access denied"}), 403
    if not is_available():
        return _ai_unavailable_response()

    data = request.get_json(silent=True) or {}
    data["_cache_feature"] = "flashcards"
    ctx, err_resp, err_code = _learniq_context(study_material, data)
    if err_resp:
        return err_resp, err_code

    from utils import ai_cache
    cached = ai_cache.get(ctx["cache_key"])
    if cached:
        return jsonify({"flashcards": cached["flashcards"], "feature": "LearnIQ", "cached": True})

    try:
        cards = learniq_flashcards(
            study_material.title, ctx["doc_text"], page_images=ctx["page_images"]
        )
        ai_cache.set(ctx["cache_key"], {"flashcards": cards})
        return jsonify({
            "flashcards": cards,
            "feature": "LearnIQ",
            "cached": False,
            "vision": ctx["use_vision"],
        })
    except ConnectionError as e:
        return jsonify({"error": "AI service is temporarily unavailable."}), 503
    except Exception as e:
        logging.exception("LearnIQ flashcards failed")
        return jsonify({"error": str(e) or "AI request failed."}), 500


@study_material_routes.route('/ai/chat/<int:course_id>', methods=['POST'])
@login_required
def learniq_chat_route(course_id):
    from utils.local_ai import learniq_chat, is_available
    from utils.ai_rate_limit import check_ai_rate_limit

    if not session.get("user_id"):
        return jsonify({"error": "Unauthorized"}), 401
    ok, retry = check_ai_rate_limit()
    if not ok:
        return jsonify({"error": f"Rate limit exceeded. Retry in {retry}s.", "retry_after": retry}), 429

    study_material = StudyMaterial.query.get_or_404(course_id)
    assert_tenant_access(study_material)
    user = User.query.get(session["user_id"])
    if not can_access_study_material(user, study_material):
        return jsonify({"error": "Access denied"}), 403
    if not is_available():
        return _ai_unavailable_response()

    data = request.get_json(silent=True) or {}
    message = (data.get("message") or "").strip()
    if not message:
        return jsonify({"error": "Message is required."}), 400

    data["_cache_feature"] = "chat"
    ctx, err_resp, err_code = _learniq_context(study_material, data)
    if err_resp:
        return err_resp, err_code

    try:
        reply = learniq_chat(
            study_material.title,
            ctx["doc_text"],
            message,
            data.get("history", []),
            page_images=ctx["page_images"],
        )
        return jsonify({"reply": reply, "feature": "LearnIQ", "vision": ctx["use_vision"]})
    except ConnectionError as e:
        return jsonify({"error": "AI service is temporarily unavailable."}), 503
    except Exception as e:
        logging.exception("LearnIQ chat failed")
        return jsonify({"error": str(e) or "AI request failed."}), 500


@study_material_routes.route('/ai/stream/chat/<int:course_id>', methods=['POST'])
@login_required
def learniq_stream_chat_route(course_id):
    from utils.local_ai import learniq_chat_stream, is_available
    from utils.ai_rate_limit import check_ai_rate_limit

    if not session.get("user_id"):
        return jsonify({"error": "Unauthorized"}), 401
    ok, retry = check_ai_rate_limit()
    if not ok:
        return jsonify({"error": f"Rate limit exceeded. Retry in {retry}s.", "retry_after": retry}), 429

    study_material = StudyMaterial.query.get_or_404(course_id)
    assert_tenant_access(study_material)
    user = User.query.get(session["user_id"])
    if not can_access_study_material(user, study_material):
        return jsonify({"error": "Access denied"}), 403
    if not is_available():
        return _ai_unavailable_response()

    data = request.get_json(silent=True) or {}
    message = (data.get("message") or "").strip()
    if not message:
        return jsonify({"error": "Message is required."}), 400

    data["_cache_feature"] = "chat"
    ctx, err_resp, err_code = _learniq_context(study_material, data)
    if err_resp:
        return err_resp, err_code

    def generate():
        try:
            for chunk in learniq_chat_stream(
                study_material.title,
                ctx["doc_text"],
                message,
                data.get("history", []),
                page_images=ctx["page_images"],
            ):
                yield f"data: {json_lib.dumps({'text': chunk})}\n\n"
            yield f"data: {json_lib.dumps({'done': True, 'vision': ctx['use_vision']})}\n\n"
        except ConnectionError as e:
            logging.error("LearnIQ chat stream failed: %s", e)
            yield f"data: {json_lib.dumps({'error': 'AI service unavailable. Check that Ollama is running.'})}\n\n"
        except Exception as e:
            logging.exception("LearnIQ chat stream error")
            yield f"data: {json_lib.dumps({'error': str(e) or 'AI request failed.'})}\n\n"

    return Response(stream_with_context(generate()), mimetype="text/event-stream")


@study_material_routes.route('/ai/stream/sample_questions/<int:course_id>', methods=['POST'])
@login_required
def learniq_stream_sample_questions_route(course_id):
    from utils.local_ai import learniq_sample_questions_stream, is_available
    from utils.ai_rate_limit import check_ai_rate_limit

    if not session.get("user_id"):
        return jsonify({"error": "Unauthorized"}), 401
    ok, retry = check_ai_rate_limit()
    if not ok:
        return jsonify({"error": f"Rate limit exceeded. Retry in {retry}s.", "retry_after": retry}), 429

    study_material = StudyMaterial.query.get_or_404(course_id)
    assert_tenant_access(study_material)
    user = User.query.get(session["user_id"])
    if not can_access_study_material(user, study_material):
        return jsonify({"error": "Access denied"}), 403
    if not is_available():
        return _ai_unavailable_response()

    data = request.get_json(silent=True) or {}
    data["_cache_feature"] = "sample_questions"
    ctx, err_resp, err_code = _learniq_context(study_material, data)
    if err_resp:
        return err_resp, err_code

    from utils import ai_cache
    cached = ai_cache.get(ctx["cache_key"])

    def generate():
        if cached:
            yield f"data: {json_lib.dumps({'text': cached['questions'], 'cached': True})}\n\n"
            yield f"data: {json_lib.dumps({'done': True, 'cached': True})}\n\n"
            return
        parts = []
        try:
            for chunk in learniq_sample_questions_stream(
                study_material.title, ctx["doc_text"], page_images=ctx["page_images"]
            ):
                parts.append(chunk)
                yield f"data: {json_lib.dumps({'text': chunk})}\n\n"
            ai_cache.set(ctx["cache_key"], {"questions": "".join(parts)})
            yield f"data: {json_lib.dumps({'done': True, 'cached': False, 'vision': ctx['use_vision']})}\n\n"
        except ConnectionError:
            yield f"data: {json_lib.dumps({'error': 'AI service is temporarily unavailable.'})}\n\n"

    return Response(stream_with_context(generate()), mimetype="text/event-stream")


@study_material_routes.route('/ai/sample_questions/<int:course_id>', methods=['POST'])
@study_material_routes.route('/ai/quiz/<int:course_id>', methods=['POST'])
@login_required
def learniq_quiz_json(course_id):
    """Non-streaming Quiz Me fallback (also used when SSE route is unavailable)."""
    from utils.local_ai import learniq_sample_questions, is_available
    from utils.ai_rate_limit import check_ai_rate_limit

    if not session.get("user_id"):
        return jsonify({"error": "Unauthorized"}), 401
    ok, retry = check_ai_rate_limit()
    if not ok:
        return jsonify({"error": f"Rate limit exceeded. Retry in {retry}s.", "retry_after": retry}), 429

    study_material = StudyMaterial.query.get_or_404(course_id)
    assert_tenant_access(study_material)
    user = User.query.get(session["user_id"])
    if not can_access_study_material(user, study_material):
        return jsonify({"error": "Access denied"}), 403
    if not is_available():
        return _ai_unavailable_response()

    data = request.get_json(silent=True) or {}
    data["_cache_feature"] = "sample_questions"
    ctx, err_resp, err_code = _learniq_context(study_material, data)
    if err_resp:
        return err_resp, err_code

    from utils import ai_cache
    cached = ai_cache.get(ctx["cache_key"])
    if cached:
        return jsonify({
            "questions": cached["questions"],
            "feature": "LearnIQ",
            "cached": True,
            "vision": ctx["use_vision"],
        })

    try:
        questions = learniq_sample_questions(
            study_material.title, ctx["doc_text"], page_images=ctx["page_images"]
        )
        ai_cache.set(ctx["cache_key"], {"questions": questions})
        return jsonify({
            "questions": questions,
            "feature": "LearnIQ",
            "cached": False,
            "vision": ctx["use_vision"],
        })
    except ConnectionError:
        return jsonify({"error": "AI service is temporarily unavailable."}), 503
    except Exception as e:
        logging.exception("LearnIQ quiz failed")
        return jsonify({"error": str(e) or "AI request failed."}), 500


@study_material_routes.route('/get_dropdowns', methods=['GET'])
@login_required
def get_dropdowns():
    """Fetch tenant-scoped Levels, Categories, and Designations."""
    try:
        levels = tenant_levels_query().order_by(Level.level_number.asc()).all()
        categories = tenant_categories_query().order_by(Category.id.asc()).all()
        designations = tenant_designations_query().order_by(Designation.id.asc()).all()

        # Constructing JSON response
        data = {
            "levels": [{"id": level.id, "number": level.level_number} for level in levels],
            "categories": [{"id": category.id, "name": category.name} for category in categories],
            "designations": [
                {
                    "id": designation.id,
                    "title": designation.title,
                    "starting_level": designation.starting_level or 1,
                }
                for designation in designations
            ],
        }

        return jsonify(data), 200
    except Exception as e:
        logging.error(f"Error fetching dropdowns: {e}")
        return jsonify({"error": "Failed to fetch dropdowns"}), 500


@study_material_routes.route('/study_materials/import-scorm', methods=['POST'])
@login_required
def import_scorm_package():
    """Import SCORM 1.2 / 2004 or xAPI zip as a course stub (admin upload)."""
    from utils.user_access import can_upload_study_materials
    from utils.scorm_import import import_scorm_to_course
    import tempfile
    import os

    if not can_upload_study_materials(current_user):
        flash("You do not have permission to import courses.", "error")
        return redirect(url_for('study_material_routes.upload_page'))

    upload = request.files.get('scorm_zip')
    if not upload or not upload.filename.lower().endswith('.zip'):
        flash("Upload a .zip SCORM or xAPI package.", "error")
        return redirect(url_for('study_material_routes.upload_page'))

    title_override = (request.form.get('title') or '').strip() or None
    with tempfile.NamedTemporaryFile(suffix='.zip', delete=False) as tmp:
        upload.save(tmp.name)
        tmp_path = tmp.name
    try:
        material, info = import_scorm_to_course(
            tmp_path,
            tenant_id=user_tenant_id(),
            title_override=title_override,
        )
        flash(f"Imported {info.format.upper()} course: {material.title}", "success")
        return redirect(url_for('study_material_routes.view_course', course_id=material.id))
    except Exception as exc:
        logging.error("SCORM import failed: %s", exc)
        flash("Could not import SCORM package.", "error")
        return redirect(url_for('study_material_routes.upload_page'))
    finally:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass

